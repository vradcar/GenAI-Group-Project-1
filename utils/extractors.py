"""
Text Extractors — file-type-specific text extraction.

Responsibilities:
  - extract_pdf(path)  → str   (via PyMuPDF)
  - extract_pptx(path) → str   (via python-pptx)
  - extract_txt(path)  → str   (plain read)
  - extract_url(url)   → str   (via trafilatura)
"""

import fitz  # PyMuPDF
import trafilatura
from pptx import Presentation


def extract_pdf(file_path: str) -> str:
    """Extract text from a PDF file using PyMuPDF (fitz)."""
    doc = fitz.open(file_path)
    pages = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            pages.append(text)
    doc.close()
    return "\n\n".join(pages).strip()


def extract_pptx(file_path: str) -> str:
    """Extract text from all slides in a PPTX file using python-pptx."""
    prs = Presentation(file_path)
    slide_texts = []
    for slide in prs.slides:
        shape_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shape_texts.append(shape.text.strip())
        if shape_texts:
            slide_texts.append("\n".join(shape_texts))
    return "\n\n".join(slide_texts)


def extract_txt(file_path: str) -> str:
    """Read and return the contents of a plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def extract_url(url: str) -> str:
    """Fetch a URL and extract the main article content using trafilatura."""
    downloaded = trafilatura.fetch_url(url)
    if not downloaded:
        raise ValueError(f"Could not fetch content from URL: {url}")
    text = trafilatura.extract(downloaded)
    if not text:
        raise ValueError(f"Could not extract readable text from URL: {url}")
    return text
